import os
import re
import requests
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from moviepy.video.VideoClip import ImageClip
from moviepy.video.io.VideoFileClip import VideoFileClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
from moviepy.video.compositing.CompositeVideoClip import CompositeVideoClip
from moviepy.video.compositing.concatenate import concatenate_videoclips
from pdf2image import convert_from_path
import pyttsx3
import subprocess

app = Flask(__name__)

# ---------------- 路径 ----------------
UPLOAD_DIR = "uploads"
OUT_DIR = "outputs"
DIGITAL_ACTOR_PATH = "static/1.gif"  # 数字人素材
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUT_DIR, exist_ok=True)

# ================= 扣子大模型 =================
COZE_API_KEY = "pat_GDsW1EkdLisaq2mGbgjbXL9vhXa3eqRsT5UV6tquTGLeANzkp7V4q9EW5p3v4R8D"
BOT_ID = "7584362394714472454"

def generate_lecture_text(ppt_text):
    """
    PPT 提纲 → 教师讲课式讲稿
    """
    prompt = f"""
你是一名高校教师，请根据下面 PPT 内容，
生成自然、口语化、适合课堂讲解的讲稿，
不要逐字朗读 PPT，要适当扩展。

PPT 内容：
{ppt_text}
"""

    url = "https://api.coze.cn/open_api/v2/chat"
    headers = {
        "Authorization": f"Bearer {COZE_API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "bot_id": BOT_ID,
        "user": "teacher",
        "query": prompt,
        "stream": False
    }

    resp = requests.post(url, headers=headers, json=payload, timeout=60)
    data = resp.json()

    for msg in data.get("messages", []):
        if msg.get("type") == "answer":
            return msg.get("content")

    return ppt_text


# ================= 拆句（核心） =================
def split_sentences(text):
    parts = re.split(r'(。|！|\!|？|\?|，|,| )', text)
    sentences = []
    buf = ""

    for p in parts:
        buf += p
        if p in "。！？!?，, ":
            s = buf.strip()
            if len(s) > 2:
                sentences.append(s)
            buf = ""

    if buf.strip():
        sentences.append(buf.strip())

    return sentences

# ---------------- PPT 转图片 ----------------
def ppt_to_images_real(ppt_path):
    SOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
    subprocess.run([
        SOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        ppt_path,
        "--outdir", OUT_DIR
    ], check=True)

    pdf_path = os.path.join(OUT_DIR, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    POPPLER_PATH = r"C:\Users\陈\poppler\poppler-25.12.0\Library\bin"

    images = convert_from_path(pdf_path, dpi=200, poppler_path=POPPLER_PATH)
    paths = []
    for i, img in enumerate(images):
        path = os.path.join(OUT_DIR, f"slide_{i}.png")
        img.save(path)
        paths.append(path)
    return paths

# ---------------- 提取 PPT 文本 ----------------
def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    texts = []
    for slide in prs.slides:
        t = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t.append(shape.text)
        texts.append("\n".join(t))
    return texts

# ---------------- TTS ----------------
def tts(text, path):
    engine = pyttsx3.init()
    engine.save_to_file(text, path)
    engine.runAndWait()

# ---------------- 字幕图（自动换行） ----------------
def subtitle_image(text, width=1280, height=150, font_size=40, line_spacing=10):
    img = Image.new("RGBA", (width, height), (0, 0, 0, 180))
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("msyh.ttc", font_size)
    except:
        font = ImageFont.load_default()

    def wrap_text(text, font, max_width):
        chars = list(text)
        lines = []
        line = ""
        for c in chars:
            line += c
            w, _ = draw.textsize(line, font=font)
            if w > max_width:
                lines.append(line[:-1])
                line = c
        if line:
            lines.append(line)
        return lines

    lines = wrap_text(text, font, width - 80)
    total_height = len(lines) * (font_size + line_spacing)
    y_start = (height - total_height) // 2
    y = y_start
    for line in lines:
        w, h = draw.textsize(line, font=font)
        x = (width - w) // 2
        draw.text((x, y), line, fill="white", font=font)
        y += font_size + line_spacing
    return img

# ---------------- 路由 ----------------
@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    if "ppt" not in request.files:
        return "未上传 PPT 文件", 400

    ppt = request.files["ppt"]
    ppt_path = os.path.join(UPLOAD_DIR, ppt.filename)
    ppt.save(ppt_path)

    # 1️⃣ PPT → 图片
    try:
        slide_imgs = ppt_to_images_real(ppt_path)
    except Exception as e:
        return f"PPT 转图片失败：{e}", 500

    # 2️⃣ 提取 PPT 文本
    slide_texts = extract_ppt_text(ppt_path)

    clips = []

    for i, img_path in enumerate(slide_imgs):
        slide_img = ImageClip(img_path)
        # 3️⃣ 拆句
        explain_text = slide_texts[i]  # 这里可以换成你调用 COZE API 的讲解文本
        sentences = [s.strip() for s in explain_text.replace("\n","").split("。") if s.strip()]
        if not sentences:
            sentences = [explain_text]

        slide_clips = []

        for j, sentence in enumerate(sentences):
            # TTS
            audio_path = os.path.join(OUT_DIR, f"a_{i}_{j}.mp3")
            tts(sentence, audio_path)
            audio = AudioFileClip(audio_path)
            dur = audio.duration

            # 数字人素材
            if DIGITAL_ACTOR_PATH.endswith(".gif"):
                actor_clip = VideoFileClip(DIGITAL_ACTOR_PATH).loop(duration=dur).resize(width=200)
            else:
                actor_clip = ImageClip(DIGITAL_ACTOR_PATH).set_duration(dur).resize(width=200)
            actor_clip = actor_clip.set_position(("right","bottom"))

            # 字幕
            sub_img = subtitle_image(sentence)
            sub_path = os.path.join(OUT_DIR, f"s_{i}_{j}.png")
            sub_img.save(sub_path)
            subtitle_clip = ImageClip(sub_path).set_duration(dur).set_position(("center","bottom"))

            # 合成
            clip = CompositeVideoClip([slide_img.set_duration(dur), actor_clip, subtitle_clip]).set_audio(audio)
            slide_clips.append(clip)

        # 拼接每页句子
        page_clip = concatenate_videoclips(slide_clips)
        clips.append(page_clip)

    if not clips:
        return jsonify({"error":"未生成视频片段"}),500

    # 拼接所有页
    final = concatenate_videoclips(clips)
    out_path = os.path.join(OUT_DIR, "result.mp4")
    final.write_videofile(out_path, fps=24)

    return send_file(out_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
